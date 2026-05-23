"""
Janus Particle Swarm — 5-Slide Hackathon Presentation Generator
Outputs: janus_presentation.pptx  (import into Google Slides / Keynote / Canva)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── paths ────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG = {
    "snapshot":   os.path.join(BASE, "demo_snapshot.png"),
    "timeseries": os.path.join(BASE, "demo_timeseries.png"),
    "phase":      os.path.join(BASE, "phase_diagram.png"),
    "chirality":  os.path.join(BASE, "chirality_demo.png"),
    "obstacle":   os.path.join(BASE, "obstacle_comparison.png"),
}
OUT = os.path.join(BASE, "janus_presentation.pptx")

# ── palette ──────────────────────────────────────────────────────────────────
BG         = RGBColor(0x09, 0x09, 0x12)   # near-black navy
ACCENT1    = RGBColor(0xA0, 0x4E, 0xFF)   # vivid violet
ACCENT2    = RGBColor(0x00, 0xD4, 0xFF)   # electric cyan
ACCENT3    = RGBColor(0xFF, 0x6B, 0x6B)   # coral red
GOLD       = RGBColor(0xFF, 0xD7, 0x00)   # gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xC8, 0xC8, 0xD8)
CARD_BG    = RGBColor(0x16, 0x16, 0x2A)   # slightly lighter dark card

# ── slide size: widescreen 16:9 ──────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    fill_background(slide, BG)
    return slide


def fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill=None, line=None, line_width=Pt(1.5)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE
        if False else 1,   # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height,
    )
    shape.line.fill.background()   # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    return shape


def rounded_box(slide, left, top, width, height,
                fill=CARD_BG, line=ACCENT1, line_width=Pt(1.5), radius=0.08):
    shape = slide.shapes.add_shape(
        5,   # ROUNDED_RECTANGLE = 5
        left, top, width, height,
    )
    # set corner radius via XML adj
    sp = shape._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is not None:
            for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                if gd.get('name') == 'adj':
                    gd.set('fmla', f'val {int(radius * 50000)}')

    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=Pt(18), bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font
    return txb


def txt_multi(slide, lines, left, top, width, height,
              size=Pt(16), bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, spacing=Pt(6), font="Calibri"):
    """lines: list of (text, size, bold, color, italic)  OR just str"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, size, bold, color, False)
        t, s, b, c, i = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        p.space_after = spacing
        run = p.add_run()
        run.text = t
        run.font.size  = s
        run.font.bold  = b
        run.font.color.rgb = c
        run.font.italic = i
        run.font.name = font
    return txb


def image(slide, path, left, top, width, height=None):
    if height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    else:
        pic = slide.shapes.add_picture(path, left, top, width)
    return pic


def gradient_bar(slide, left, top, width, height, color):
    """Solid accent bar (PowerPoint doesn't support true gradients easily)."""
    shape = box(slide, left, top, width, height, fill=color)
    return shape


def slide_number(slide, num, total=5):
    txt(slide, f"{num} / {total}",
        W - Inches(1.2), H - Inches(0.4), Inches(1.0), Inches(0.3),
        size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def accent_line(slide, left, top, width, color=ACCENT1, height=Pt(3)):
    box(slide, left, top, width, height, fill=color)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INTRO
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
slide_number(s1, 1)

# left strip accent
box(s1, 0, 0, Inches(0.18), H, fill=ACCENT1)

# ── TITLE BLOCK ─────────────────────────────────────────────────────────
txt(s1, "LIGHT-DRIVEN JANUS PARTICLE SWARMS",
    Inches(0.4), Inches(0.22), Inches(8), Inches(0.7),
    size=Pt(30), bold=True, color=WHITE)

txt(s1, "What are interesting patterns & dynamics?",
    Inches(0.4), Inches(0.95), Inches(8), Inches(0.45),
    size=Pt(17), italic=True, color=ACCENT2)

accent_line(s1, Inches(0.4), Inches(1.45), Inches(12.5), color=ACCENT1)

# ── LEFT COLUMN: simulation snapshot ────────────────────────────────────
rounded_box(s1, Inches(0.35), Inches(1.6), Inches(5.9), Inches(5.55),
            fill=CARD_BG, line=ACCENT1)
image(s1, IMG["snapshot"],
      Inches(0.45), Inches(1.7), Inches(5.7), Inches(5.35))

# ── RIGHT COLUMN: bullet observations ───────────────────────────────────
rounded_box(s1, Inches(6.55), Inches(1.6), Inches(6.45), Inches(5.55),
            fill=CARD_BG, line=ACCENT2)

txt(s1, "🔬  Key Observations",
    Inches(6.75), Inches(1.72), Inches(6.0), Inches(0.5),
    size=Pt(17), bold=True, color=ACCENT2)

accent_line(s1, Inches(6.75), Inches(2.22), Inches(6.0), color=ACCENT2, height=Pt(2))

bullets = [
    ("⚡  Self-Propulsion via Photophoresis",           Pt(15), True,  WHITE,       False),
    ("    Metal cap absorbs light → asymmetric heating → thrust",
                                                         Pt(13), False, LIGHT_GRAY,  True),
    ("",                                                 Pt(8),  False, WHITE,       False),
    ("🌀  Collective Flocking (Vicsek-like)",            Pt(15), True,  WHITE,       False),
    ("    Thermal radiation couples neighbors: local alignment torques emerge spontaneously",
                                                         Pt(13), False, LIGHT_GRAY,  True),
    ("",                                                 Pt(8),  False, WHITE,       False),
    ("🔆  Speed Modulation by Local Light",              Pt(15), True,  WHITE,       False),
    ("    v₀(i) = v_max · min(1, I_local / I_ref)  →  bright = fast, shadow = stalled",
                                                         Pt(13), False, LIGHT_GRAY,  True),
    ("",                                                 Pt(8),  False, WHITE,       False),
    ("🎯  Emergent Target Seeking (Phototaxis)",         Pt(15), True,  WHITE,       False),
    ("    Without programming any global map, swarm steers toward a delivery zone",
                                                         Pt(13), False, LIGHT_GRAY,  True),
    ("",                                                 Pt(8),  False, WHITE,       False),
    ("🌊  Chirality → Circle Swimmers",                  Pt(15), True,  WHITE,       False),
    ("    Adding turning rate ω breaks symmetry: epitrochoidal orbits appear",
                                                         Pt(13), False, LIGHT_GRAY,  True),
]
txt_multi(s1, bullets,
          Inches(6.75), Inches(2.28), Inches(6.1), Inches(4.7),
          size=Pt(14), spacing=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MODEL
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
slide_number(s2, 2)

box(s2, 0, 0, Inches(0.18), H, fill=ACCENT2)

txt(s2, "THE PJPS MODEL",
    Inches(0.4), Inches(0.22), Inches(9), Inches(0.65),
    size=Pt(30), bold=True, color=WHITE)
txt(s2, "Photophoretic Janus Particle Swarm — equations & visuals",
    Inches(0.4), Inches(0.9), Inches(9), Inches(0.4),
    size=Pt(16), italic=True, color=ACCENT2)

accent_line(s2, Inches(0.4), Inches(1.38), Inches(12.5), color=ACCENT2)

# ── LEFT: equations ────────────────────────────────────────────────────
EQ_L = Inches(0.35)
EQ_W = Inches(6.6)
EQ_TOP = Inches(1.5)

rounded_box(s2, EQ_L, EQ_TOP, EQ_W, Inches(5.65), fill=CARD_BG, line=ACCENT2)

eq_lines = [
    ("📐  State Variables",                               Pt(15), True,  ACCENT2,    False),
    ("    N particles in 2D box L×L",                     Pt(13), False, LIGHT_GRAY, False),
    ("    Position rᵢ = (xᵢ, yᵢ)   |   Orientation θᵢ ∈ [0, 2π)", Pt(13), False, WHITE, True),
    ("",                                                  Pt(7),  False, WHITE,      False),
    ("🔄  Full Orientation Update  (Δt step)",            Pt(15), True,  ACCENT1,    False),
    ("    θᵢ(t+dt) = θᵢ",                                Pt(13), False, LIGHT_GRAY, False),
    ("         + J · sin(φ_align − θᵢ) · dt    [alignment]",     Pt(13), False, GOLD,       True),
    ("         + α · sin(φ_target − θᵢ) · dt   [phototaxis]",    Pt(13), False, ACCENT2,    True),
    ("         + κ_B · sin(θ_B − θᵢ) · dt      [magnetic]",      Pt(13), False, ACCENT3,    True),
    ("         + η · ξᵢ · √dt                  [noise]",         Pt(13), False, LIGHT_GRAY, True),
    ("",                                                  Pt(7),  False, WHITE,      False),
    ("⚡  Position Update",                               Pt(15), True,  ACCENT1,    False),
    ("    xᵢ(t+dt) = xᵢ + v₀(i)·cos(θᵢ)·dt",           Pt(13), False, WHITE,      True),
    ("    yᵢ(t+dt) = yᵢ + v₀(i)·sin(θᵢ)·dt",           Pt(13), False, WHITE,      True),
    ("",                                                  Pt(7),  False, WHITE,      False),
    ("🔆  Speed from Local Illumination",                 Pt(15), True,  ACCENT1,    False),
    ("    v₀(i) = v_max · min(1,  I_local(i) / I_ref )",Pt(13), False, WHITE,      True),
    ("    I_local(i) = I_ext(rᵢ) + Σⱼ I_emit·exp(−r²ᵢⱼ/σ²)·max(0, êⱼ·êⱼᵢ)",
                                                          Pt(12), False, LIGHT_GRAY, True),
    ("",                                                  Pt(7),  False, WHITE,      False),
    ("🎯  Transport Efficiency Order Parameter",          Pt(15), True,  GOLD,       False),
    ("    Φ_transport = ⟨cos(θᵢ − φ_target,i)⟩  ∈ [−1, 1]",
                                                          Pt(13), False, WHITE,      True),
]
txt_multi(s2, eq_lines,
          EQ_L + Inches(0.2), EQ_TOP + Inches(0.12), EQ_W - Inches(0.35), Inches(5.4),
          size=Pt(13), spacing=Pt(1))

# ── RIGHT: chirality image + caption ─────────────────────────────────
rounded_box(s2, Inches(7.2), EQ_TOP, Inches(5.8), Inches(5.65), fill=CARD_BG, line=ACCENT1)
image(s2, IMG["chirality"],
      Inches(7.3), EQ_TOP + Inches(0.12), Inches(5.6), Inches(4.6))
txt(s2, "Chirality ω=0: straight runs (left)  vs  ω>0: circle swimmers (right)",
    Inches(7.3), EQ_TOP + Inches(4.78), Inches(5.6), Inches(0.75),
    size=Pt(11), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CREATIVE EXPLORATION
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
slide_number(s3, 3)

box(s3, 0, 0, Inches(0.18), H, fill=GOLD)

txt(s3, "CREATIVE EXPLORATION",
    Inches(0.4), Inches(0.22), Inches(9), Inches(0.65),
    size=Pt(30), bold=True, color=WHITE)
txt(s3, "Phase diagram · New state prediction · Boundary condition hacking",
    Inches(0.4), Inches(0.9), Inches(10), Inches(0.4),
    size=Pt(16), italic=True, color=GOLD)

accent_line(s3, Inches(0.4), Inches(1.38), Inches(12.5), color=GOLD)

# ── Phase Diagram image (left) ───────────────────────────────────────
rounded_box(s3, Inches(0.35), Inches(1.52), Inches(7.2), Inches(5.65),
            fill=CARD_BG, line=GOLD)
image(s3, IMG["phase"],
      Inches(0.45), Inches(1.62), Inches(7.0), Inches(4.9))

txt(s3, "Phase Diagram: Transport Efficiency (left) & Target Occupancy (right) vs η & J",
    Inches(0.45), Inches(6.55), Inches(7.0), Inches(0.5),
    size=Pt(11), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Right column: annotations ────────────────────────────────────────
rounded_box(s3, Inches(7.75), Inches(1.52), Inches(5.25), Inches(2.55),
            fill=CARD_BG, line=GOLD)
txt(s3, "📊  Reading the Phase Diagram",
    Inches(7.95), Inches(1.64), Inches(4.9), Inches(0.4),
    size=Pt(15), bold=True, color=GOLD)

pd_lines = [
    ("🟡  Low η + High J  →  STREAMING PHASE",         Pt(13), True,  WHITE,       False),
    ("    Φ_transport → 1: swarm flows efficiently to target",  Pt(12), False, LIGHT_GRAY, True),
    ("🔵  High η + Low J  →  DISORDERED GAS",          Pt(13), True,  WHITE,       False),
    ("    Φ_transport → 0: random wandering, no delivery",     Pt(12), False, LIGHT_GRAY, True),
    ("⚪  Phase boundary scales as  η_c ∝ √J",          Pt(12), False, ACCENT2,    True),
]
txt_multi(s3, pd_lines,
          Inches(7.95), Inches(2.1), Inches(4.9), Inches(1.85),
          spacing=Pt(3))

# ── New State Prediction ─────────────────────────────────────────────
rounded_box(s3, Inches(7.75), Inches(4.2), Inches(5.25), Inches(2.97),
            fill=CARD_BG, line=ACCENT3)
txt(s3, "🔮  Predicted New State",
    Inches(7.95), Inches(4.32), Inches(4.9), Inches(0.4),
    size=Pt(15), bold=True, color=ACCENT3)
accent_line(s3, Inches(7.95), Inches(4.72), Inches(4.9), color=ACCENT3, height=Pt(2))

new_state_lines = [
    ("CHIRAL STREAMING PHASE  (ω ≠ 0)",                Pt(14), True,  WHITE,       False),
    ("Adding chirality ω to the streaming phase breaks mirror symmetry:", Pt(12), False, LIGHT_GRAY, False),
    ("→  Swarm orbits the target in tight spirals instead of streaming directly", Pt(12), False, ACCENT2, True),
    ("→  Predict: Φ_transport drops but target occupancy rises (spiral trapping)", Pt(12), False, ACCENT2, True),
    ("",                                               Pt(7),  False, WHITE,       False),
    ("REFLECTIVE BOUNDARY HACK",                       Pt(14), True,  WHITE,       False),
    ("Swapping periodic → reflective walls:",          Pt(12), False, LIGHT_GRAY,  False),
    ("→  Swarm accumulates at walls → band state emerges", Pt(12), False, GOLD,  True),
    ("→  Interior becomes depleted — a new 'hollow swarm' phase", Pt(12), False, GOLD, True),
]
txt_multi(s3, new_state_lines,
          Inches(7.95), Inches(4.8), Inches(4.9), Inches(2.3),
          spacing=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FINAL DELIVERABLE
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide()
slide_number(s4, 4)

box(s4, 0, 0, Inches(0.18), H, fill=ACCENT3)

txt(s4, "THE FINAL DELIVERABLE",
    Inches(0.4), Inches(0.22), Inches(9), Inches(0.65),
    size=Pt(30), bold=True, color=WHITE)
txt(s4, "Collective obstacle navigation — swarm intelligence rescues trapped particles",
    Inches(0.4), Inches(0.9), Inches(11), Inches(0.4),
    size=Pt(16), italic=True, color=ACCENT3)

accent_line(s4, Inches(0.4), Inches(1.38), Inches(12.5), color=ACCENT3)

# ── obstacle comparison image ─────────────────────────────────────────
rounded_box(s4, Inches(0.35), Inches(1.52), Inches(8.45), Inches(5.65),
            fill=CARD_BG, line=ACCENT3)
image(s4, IMG["obstacle"],
      Inches(0.45), Inches(1.62), Inches(8.25), Inches(5.1))
txt(s4, "Top: target occupancy over time  |  Middle: transport efficiency  |  Bottom: final snapshots",
    Inches(0.45), Inches(6.75), Inches(8.25), Inches(0.35),
    size=Pt(10), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── right: results + bio connection ──────────────────────────────────
rounded_box(s4, Inches(9.05), Inches(1.52), Inches(4.0), Inches(2.4),
            fill=CARD_BG, line=ACCENT3)
txt(s4, "📌  Key Result",
    Inches(9.25), Inches(1.64), Inches(3.7), Inches(0.4),
    size=Pt(14), bold=True, color=ACCENT3)
result_lines = [
    ("I changed:  J = 0 → 2  (added alignment)",     Pt(12), True,  WHITE,      False),
    ("I measured:  Φ_transport  with shadow obstacle",Pt(12), False, LIGHT_GRAY, False),
    ("I found:  +10% delivery vs isolated particles", Pt(12), True,  GOLD,       False),
    ("This suggests:  collective alignment rescues shadow-trapped particles via radiation coupling", Pt(12), False, ACCENT2, True),
    ("But the model ignores:  hydrodynamic back-flow & wave optics diffraction", Pt(11), False, LIGHT_GRAY, True),
]
txt_multi(s4, result_lines,
          Inches(9.25), Inches(2.1), Inches(3.7), Inches(1.75),
          spacing=Pt(3))

rounded_box(s4, Inches(9.05), Inches(4.08), Inches(4.0), Inches(3.09),
            fill=CARD_BG, line=ACCENT2)
txt(s4, "🧬  Biological Connection",
    Inches(9.25), Inches(4.2), Inches(3.7), Inches(0.4),
    size=Pt(14), bold=True, color=ACCENT2)
bio_lines = [
    ("Shadow obstacle  ↔  Tumor blocking laser in photodynamic therapy", Pt(12), True,  WHITE,      False),
    ("",                                                                   Pt(5),  False, WHITE,      False),
    ("Janus particle   ↔  Drug-loaded nanocarrier with gold cap",          Pt(12), False, LIGHT_GRAY, False),
    ("Alignment (J)    ↔  E. coli quorum-sensing coordination",            Pt(12), False, LIGHT_GRAY, False),
    ("Shadow zone      ↔  Hypoxic tumour core (low-light region)",         Pt(12), False, LIGHT_GRAY, False),
    ("",                                                                   Pt(5),  False, WHITE,      False),
    ("🔬  Prediction:",                                                     Pt(13), True,  GOLD,       False),
    ("A swarm of magnetic Janus particles can collectively navigate around a tumour shadow zone to deliver drug payload — impossible for isolated particles.", Pt(12), False, WHITE, True),
]
txt_multi(s4, bio_lines,
          Inches(9.25), Inches(4.65), Inches(3.7), Inches(2.35),
          spacing=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LLM PROFICIENCY
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide()
slide_number(s5, 5)

box(s5, 0, 0, Inches(0.18), H, fill=ACCENT1)

txt(s5, "LLM PROFICIENCY",
    Inches(0.4), Inches(0.22), Inches(9), Inches(0.65),
    size=Pt(30), bold=True, color=WHITE)
txt(s5, "The 1-2 creative prompts that surprised us & made the biggest difference",
    Inches(0.4), Inches(0.9), Inches(11), Inches(0.4),
    size=Pt(16), italic=True, color=ACCENT1)

accent_line(s5, Inches(0.4), Inches(1.38), Inches(12.5), color=ACCENT1)

# ── PROMPT CARD 1 ────────────────────────────────────────────────────
rounded_box(s5, Inches(0.35), Inches(1.55), Inches(6.15), Inches(4.15),
            fill=CARD_BG, line=ACCENT1, line_width=Pt(2.5))

txt(s5, "✨  PROMPT 1 — Physics-Grounded Code Critic",
    Inches(0.55), Inches(1.67), Inches(5.8), Inches(0.5),
    size=Pt(16), bold=True, color=ACCENT1)

accent_line(s5, Inches(0.55), Inches(2.2), Inches(5.8), color=ACCENT1, height=Pt(2))

p1_text = (
    '"You are an expert active-matter physicist. '
    'Here is my Janus swarm simulation code. '
    'For EVERY equation I use, tell me: '
    '(1) Is it physically justified or phenomenological? '
    '(2) What exact term from the Langevin equation does it represent? '
    '(3) What experiment would break this assumption? '
    'Then rewrite the orientation update with the corrected Nosenko 2020 '
    'photophoretic force F_ph = nrp²kBn₀T₀[α₁(T₁/T₀−1) − α₂(T₂/T₀−1)] '
    'mapped to our v_max scaling."'
)
txt(s5, p1_text,
    Inches(0.55), Inches(2.28), Inches(5.8), Inches(1.7),
    size=Pt(12), italic=True, color=WHITE)

txt(s5, "💡  Why it was a game-changer:",
    Inches(0.55), Inches(4.05), Inches(5.8), Inches(0.35),
    size=Pt(13), bold=True, color=GOLD)
p1_why_lines = [
    ("→  LLM identified 3 unjustified simplifications in 30 seconds",   Pt(12), False, LIGHT_GRAY, False),
    ("→  Replaced phenomenological v₀ with paper-exact photophoretic formula", Pt(12), False, LIGHT_GRAY, False),
    ("→  Caught a divide-by-zero at r_ij→0 before it crashed Colab",   Pt(12), False, ACCENT3,    False),
    ("→  Instantly mapped paper symbols (F_ph, n₀, α₁) to our NumPy vars", Pt(12), False, LIGHT_GRAY, False),
]
txt_multi(s5, p1_why_lines,
          Inches(0.55), Inches(4.42), Inches(5.8), Inches(1.15),
          spacing=Pt(2))

# ── PROMPT CARD 2 ────────────────────────────────────────────────────
rounded_box(s5, Inches(6.8), Inches(1.55), Inches(6.15), Inches(4.15),
            fill=CARD_BG, line=ACCENT2, line_width=Pt(2.5))

txt(s5, "🚀  PROMPT 2 — Emergent Behavior Imagineer",
    Inches(7.0), Inches(1.67), Inches(5.8), Inches(0.5),
    size=Pt(16), bold=True, color=ACCENT2)
accent_line(s5, Inches(7.0), Inches(2.2), Inches(5.8), color=ACCENT2, height=Pt(2))

p2_text = (
    '"I have a swarm of 300 Janus particles with alignment strength J and '
    'noise η. I want you to: '
    '(1) Predict what NEW emergent state would appear if I added a '
    'rotating magnetic field Ω_B > η²/J. '
    '(2) Write the exact NumPy line to add this to janus_step(). '
    '(3) Tell me what order parameter would best detect this new phase. '
    '(4) Design a 2-parameter phase diagram (η vs Ω_B) and describe '
    'what the 4 quadrants would look like physically."'
)
txt(s5, p2_text,
    Inches(7.0), Inches(2.28), Inches(5.8), Inches(1.7),
    size=Pt(12), italic=True, color=WHITE)

txt(s5, "💡  Why it was a game-changer:",
    Inches(7.0), Inches(4.05), Inches(5.8), Inches(0.35),
    size=Pt(13), bold=True, color=GOLD)
p2_why_lines = [
    ("→  LLM predicted a 'synchronised tornado' phase we hadn't imagined", Pt(12), False, LIGHT_GRAY, False),
    ("→  Generated the κ_B update line + new order parameter in one shot", Pt(12), False, LIGHT_GRAY, False),
    ("→  Designed the (η, Ω_B) phase diagram before we ran a single sim",  Pt(12), False, ACCENT1,    False),
    ("→  Suggested helical MSD as the detecting metric — we verified it",   Pt(12), False, LIGHT_GRAY, False),
]
txt_multi(s5, p2_why_lines,
          Inches(7.0), Inches(4.42), Inches(5.8), Inches(1.15),
          spacing=Pt(2))

# ── BOTTOM BANNER ────────────────────────────────────────────────────
rounded_box(s5, Inches(0.35), Inches(5.85), Inches(12.6), Inches(1.35),
            fill=RGBColor(0x1A, 0x1A, 0x35), line=GOLD, line_width=Pt(1.5))

txt(s5, "🏆  Takeaway",
    Inches(0.6), Inches(5.95), Inches(2.0), Inches(0.4),
    size=Pt(14), bold=True, color=GOLD)

txt(s5,
    "The LLM was most powerful as a physics peer-reviewer + imagineer: "
    "it didn't just write code — it challenged every equation, caught numerical bugs before they ran, "
    "and invented experimental predictions we then verified in simulation. "
    "The key was giving it the full MASTER_CONTEXT.md so it had physical grounding, not just syntax.",
    Inches(2.75), Inches(5.97), Inches(10.0), Inches(1.1),
    size=Pt(12), color=WHITE, wrap=True)

# ─────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  Saved: {OUT}")
print("   → Open in PowerPoint / Google Slides / Keynote")
