"""
Project Hail Mary Themed - 5-Slide Presentation Generator
Outputs: janus_presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.enum.shapes

# ── Paths ────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG = {
    "snapshot":   os.path.join(BASE, "demo_snapshot.png"),
    "timeseries": os.path.join(BASE, "demo_timeseries.png"),
    "phase":      os.path.join(BASE, "phase_diagram.png"),
    "chirality":  os.path.join(BASE, "chirality_demo.png"),
    "obstacle":   os.path.join(BASE, "obstacle_comparison.png"),
}
OUT = os.path.join(BASE, "janus_presentation.pptx")

# ── Project Hail Mary Palette ────────────────────────────────────────────────
# Dark Cosmic Space background, Glowing Astrophage Gold, Spacecraft Teal, Flare Red
BG         = RGBColor(0x06, 0x0A, 0x15)   # Deep space vacuum (dark navy-black)
GOLD       = RGBColor(0xFF, 0x9D, 0x00)   # Glowing Astrophage Gold (solar absorption)
TEAL       = RGBColor(0x00, 0xE6, 0xBD)   # Spacecraft HUD Teal
CRIMSON    = RGBColor(0xFF, 0x4B, 0x4B)   # Solar Flare / Warning Red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBC, 0xC4, 0xD0)
CARD_BG    = RGBColor(0x11, 0x17, 0x26)   # Spaceship console panel color

# ── Widescreen dimensions 16:9 ──────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # Completely blank layout

# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    fill_background(slide, BG)
    return slide

def fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill=None, line=None, line_width=Pt(1.5)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()   # no line by default
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    return shape

def rounded_box(slide, left, top, width, height, fill=CARD_BG, line=GOLD, line_width=Pt(1.5), radius=0.08):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    # Adjust corner radius via XML
    sp = shape._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is not None:
            for gd in avLst.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                if gd.get('name') == 'adj':
                    gd.set('fmla', f'val {int(radius * 50000)}')

    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=Pt(18), bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font="Trebuchet MS"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font
    return txb

def txt_multi(slide, lines, left, top, width, height,
              size=Pt(14), bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, spacing=Pt(5), font="Trebuchet MS"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, size, bold, color, False)
        t, s, b, c, i = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        p.space_after = spacing
        run = p.add_run()
        run.text = t
        run.font.size  = s
        run.font.bold  = b
        run.font.color.rgb = c
        run.font.italic = i
        run.font.name = font
    return txb

def image(slide, path, left, top, width, height=None):
    if height:
        return slide.shapes.add_picture(path, left, top, width, height)
    else:
        return slide.shapes.add_picture(path, left, top, width)

def slide_number(slide, num, total=5):
    txt(slide, f"PHM MISSION LOG: SLIDE {num} OF {total}",
        W - Inches(3.5), H - Inches(0.4), Inches(3.2), Inches(0.3),
        size=Pt(10), color=TEAL, align=PP_ALIGN.RIGHT)

def accent_line(slide, left, top, width, color=GOLD, height=Pt(2)):
    box(slide, left, top, width, height, fill=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — INTRO
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
slide_number(s1, 1)

# Left glowing gold accent strip
box(s1, 0, 0, Inches(0.18), H, fill=GOLD)

# Main Title block
txt(s1, "PROJECT HAIL MARY: THE ASTROPHAGE SWARM",
    Inches(0.4), Inches(0.18), Inches(10), Inches(0.55),
    size=Pt(28), bold=True, color=WHITE)

txt(s1, "Emergent Light-Driven Swarms & Collective Phototaxis in Space",
    Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
    size=Pt(16), italic=True, color=TEAL)

txt(s1, "Demo at https://prs-016.github.io/active-matter-microbialspacesims/",
    Inches(0.4), Inches(1.08), Inches(10), Inches(0.35),
    size=Pt(14), bold=True, color=GOLD)

accent_line(s1, Inches(0.4), Inches(1.48), Inches(12.5), color=GOLD)

# Left: Cool image snapshot
rounded_box(s1, Inches(0.35), Inches(1.58), Inches(6.0), Inches(5.55), fill=CARD_BG, line=GOLD)
image(s1, IMG["snapshot"], Inches(0.45), Inches(1.68), Inches(5.8), Inches(5.35))

# Right: Observations
rounded_box(s1, Inches(6.6), Inches(1.58), Inches(6.4), Inches(5.55), fill=CARD_BG, line=TEAL)

txt(s1, "🛰️  Mission Observations",
    Inches(6.8), Inches(1.72), Inches(6.0), Inches(0.5),
    size=Pt(17), bold=True, color=TEAL)

accent_line(s1, Inches(6.8), Inches(2.2), Inches(6.0), color=TEAL, height=Pt(1.5))

bullets_s1 = [
    ("🔆  Sun-Powered Thrust (Photophoresis)",            Pt(14), True,  WHITE,   False),
    ("    Cap-specific radiation absorption creates temperature differences, driving thrust.", Pt(12), False, LIGHT_GRAY, True),
    ("",                                                 Pt(6),  False, WHITE,   False),
    ("🌀  Astrophage Collective Flocking",                Pt(14), True,  WHITE,   False),
    ("    Thermoosmotic interactions couple neighbors, triggering spontaneous group alignment.", Pt(12), False, LIGHT_GRAY, True),
    ("",                                                 Pt(6),  False, WHITE,   False),
    ("🎯  Solar Corona Targeting",                        Pt(14), True,  WHITE,   False),
    ("    Individual phototactic steering guides the swarm directly to high-intensity targets.", Pt(12), False, LIGHT_GRAY, True),
    ("",                                                 Pt(6),  False, WHITE,   False),
    ("🌊  Chiral Symmetry Breaking (Circle Swimmers)",   Pt(14), True,  WHITE,   False),
    ("    Asymmetric particle caps yield internal torque, converting straight paths into orbits.", Pt(12), False, LIGHT_GRAY, True),
]
txt_multi(s1, bullets_s1, Inches(6.8), Inches(2.28), Inches(6.0), Inches(4.7), spacing=Pt(1.5))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MODEL SPECIFICATION
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
slide_number(s2, 2)

box(s2, 0, 0, Inches(0.18), H, fill=TEAL)

txt(s2, "ASTROPHAGE MATH: THE 3D/2D LANGEVIN MODEL",
    Inches(0.4), Inches(0.22), Inches(10), Inches(0.65),
    size=Pt(28), bold=True, color=WHITE)
txt(s2, "Coupled active-matter equations governing phototactic swarm mechanics",
    Inches(0.4), Inches(0.88), Inches(10), Inches(0.4),
    size=Pt(16), italic=True, color=TEAL)

accent_line(s2, Inches(0.4), Inches(1.36), Inches(12.5), color=TEAL)

# Left: Equations box
rounded_box(s2, Inches(0.35), Inches(1.5), Inches(6.7), Inches(5.65), fill=CARD_BG, line=TEAL)

bullets_s2 = [
    ("📐  Model Coordinates",                                Pt(14), True,  TEAL,       False),
    ("    N active particles inside a periodic box L×L.",  Pt(12), False, LIGHT_GRAY, False),
    ("    State: Position rᵢ = (xᵢ, yᵢ)   |   Heading vector e_θi = (cos θᵢ, sin θᵢ)", Pt(12), False, WHITE, True),
    ("",                                                   Pt(5),  False, WHITE,      False),
    ("🔄  Swarm Orientation Langevin Dynamics",              Pt(14), True,  GOLD,       False),
    ("    θᵢ(t+dt) = θᵢ(t)",                                Pt(12), False, LIGHT_GRAY, False),
    ("         + J · sin(φ_align − θᵢ) · dt       [Vicsek Alignment]",Pt(12), False, GOLD,       True),
    ("         + α · sin(φ_target − θᵢ) · dt      [Solar Phototaxis]",Pt(12), False, TEAL,       True),
    ("         + κ_B · sin(θ_B − θᵢ) · dt         [Magnetic Gradients]",Pt(12), False, CRIMSON,    True),
    ("         + η · ξᵢ · √dt                     [Rotational Noise (Dr)]",Pt(12), False, LIGHT_GRAY, True),
    ("",                                                   Pt(5),  False, WHITE,      False),
    ("⚡  Self-Propelling Update",                           Pt(14), True,  GOLD,       False),
    ("    rᵢ(t+dt) = rᵢ(t) + v₀(i) · e_θi · dt",            Pt(12), False, WHITE,      True),
    ("",                                                   Pt(5),  False, WHITE,      False),
    ("🔆  Local Radiation Speed Modulation",                 Pt(14), True,  GOLD,       False),
    ("    v₀(i) = v_max · min(1,  I_local(i) / I_ref )",    Pt(12), False, WHITE,      True),
    ("    I_local(i) = I_ext(rᵢ) + Σⱼ I_emit · e^(−r²ᵢⱼ/σ²) · max(0, e_θj·e_ji)",Pt(11), False, LIGHT_GRAY, True),
    ("",                                                   Pt(5),  False, WHITE,      False),
    ("📊  Emergent Flow Order Parameter",                   Pt(14), True,  TEAL,       False),
    ("    Φ_transport = ⟨ cos(θᵢ − φ_target,i) ⟩ ∈ [−1, 1]  (efficiency of delivery)", Pt(12), False, WHITE, True),
]
txt_multi(s2, bullets_s2, Inches(0.55), Inches(1.62), Inches(6.3), Inches(5.4), spacing=Pt(1.2))

# Right: Chirality visual
rounded_box(s2, Inches(7.2), Inches(1.5), Inches(5.8), Inches(5.65), fill=CARD_BG, line=GOLD)
image(s2, IMG["chirality"], Inches(7.3), Inches(1.62), Inches(5.6), Inches(4.5))
txt(s2, "Chiral swimmer asymmetry breaks straight runs: ω=0 (left) vs ω>0 (right) forming stable orbits.",
    Inches(7.3), Inches(6.2), Inches(5.6), Inches(0.8),
    size=Pt(11), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CREATIVE EXPLORATION
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
slide_number(s3, 3)

box(s3, 0, 0, Inches(0.18), H, fill=GOLD)

txt(s3, "CREATIVE SWEEP: THE ASTROPHAGE PHASE DIAGRAM",
    Inches(0.4), Inches(0.22), Inches(10), Inches(0.65),
    size=Pt(28), bold=True, color=WHITE)
txt(s3, "Phase space sweeps, boundary modifications, and new state predictions",
    Inches(0.4), Inches(0.88), Inches(10), Inches(0.4),
    size=Pt(16), italic=True, color=GOLD)

accent_line(s3, Inches(0.4), Inches(1.36), Inches(12.5), color=GOLD)

# Left: Phase Diagram Image
rounded_box(s3, Inches(0.35), Inches(1.5), Inches(7.1), Inches(5.65), fill=CARD_BG, line=GOLD)
image(s3, IMG["phase"], Inches(0.45), Inches(1.6), Inches(6.9), Inches(4.9))
txt(s3, "Phase Diagram: Transport efficiency (left) and target occupancy (right) vs noise η and coupling J",
    Inches(0.45), Inches(6.55), Inches(6.9), Inches(0.55),
    size=Pt(11), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Right top: Phase breakdown
rounded_box(s3, Inches(7.65), Inches(1.5), Inches(5.35), Inches(2.65), fill=CARD_BG, line=GOLD)
txt(s3, "📊  Phase Transition Boundaries",
    Inches(7.85), Inches(1.6), Inches(5.0), Inches(0.4),
    size=Pt(16), bold=True, color=GOLD)

pd_lines = [
    ("🟡  Low Noise + Strong J  →  COHERENT FLOCKING", Pt(13), True,  WHITE,      False),
    ("    Φ_transport → 1.0: Swarm flows efficiently to solar sinks.", Pt(12), False, LIGHT_GRAY,  True),
    ("🔵  High Noise + Weak J  →  DISORDERED GAS",    Pt(13), True,  WHITE,      False),
    ("    Φ_transport → 0.0: Individual particles brownian-wander.", Pt(12), False, LIGHT_GRAY,  True),
    ("📈  Critical Boundary: Scaling law fits η_c ∝ √J", Pt(12), True,  TEAL,      True),
]
txt_multi(s3, pd_lines, Inches(7.85), Inches(2.05), Inches(5.0), Inches(2.0), spacing=Pt(2))

# Right bottom: Predictions
rounded_box(s3, Inches(7.65), Inches(4.3), Inches(5.35), Inches(2.85), fill=CARD_BG, line=TEAL)
txt(s3, "🔮  Space boundary & State Hacking",
    Inches(7.85), Inches(4.4), Inches(5.0), Inches(0.4),
    size=Pt(16), bold=True, color=TEAL)

new_state_lines = [
    ("🔴  CHIRAL TORNADO STATE (ω ≠ 0)",              Pt(13), True,  WHITE,   False),
    ("    Adding intrinsic spin breaks symmetry. Swarm spirals into a tight orbit around the solar core. Occupancy peaks as transfer efficiency falls.", Pt(11), False, LIGHT_GRAY, False),
    ("",                                               Pt(3),  False, WHITE,   False),
    ("🟢  NEUMANN WALL BOUNDARY ACCUMULATION",         Pt(13), True,  WHITE,   False),
    ("    Swapping periodic for reflective boundaries forces particles to bunch at edges, creating a high-density ring ('hollow swarm' phase).", Pt(11), False, LIGHT_GRAY, False),
]
txt_multi(s3, new_state_lines, Inches(7.85), Inches(4.8), Inches(5.0), Inches(2.3), spacing=Pt(1.5))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE FINAL DELIVERABLE (OBSTACLE NAVIGATION)
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide()
slide_number(s4, 4)

box(s4, 0, 0, Inches(0.18), H, fill=CRIMSON)

txt(s4, "ASTEROID SHADOW BLOCKAGE & SWARM RESCUE",
    Inches(0.4), Inches(0.22), Inches(10), Inches(0.65),
    size=Pt(28), bold=True, color=WHITE)
txt(s4, "Cooperative flocking coordinates bypass around solar eclipse zones",
    Inches(0.4), Inches(0.88), Inches(10), Inches(0.4),
    size=Pt(16), italic=True, color=CRIMSON)

accent_line(s4, Inches(0.4), Inches(1.36), Inches(12.5), color=CRIMSON)

# Left: Obstacle Comparison Image
rounded_box(s4, Inches(0.35), Inches(1.5), Inches(8.35), Inches(5.65), fill=CARD_BG, line=CRIMSON)
image(s4, IMG["obstacle"], Inches(0.45), Inches(1.6), Inches(8.15), Inches(5.1))
txt(s4, "Top: Target occupancy | Middle: Transport efficiency | Bottom: Swarm avoiding shadow blockage.",
    Inches(0.45), Inches(6.75), Inches(8.15), Inches(0.35),
    size=Pt(10), italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Right top: Key Result
rounded_box(s4, Inches(8.9), Inches(1.5), Inches(4.1), Inches(2.55), fill=CARD_BG, line=CRIMSON)
txt(s4, "🎯  Key Sweeping Metric",
    Inches(9.1), Inches(1.6), Inches(3.7), Inches(0.4),
    size=Pt(16), bold=True, color=CRIMSON)

res_lines = [
    ("Control Variable: Swarmed (J=2) vs Free (J=0)",  Pt(12), True,  WHITE,      False),
    ("Measured metric: Shadow bypass rate (Φ)",         Pt(12), False, LIGHT_GRAY, False),
    ("Key Discovery: +10% delivery speed",             Pt(12), True,  GOLD,       False),
    ("Mechanism: Neighbor-alignment acts as collective memory, dragging stalled shadow-particles back to light.", Pt(11), False, LIGHT_GRAY, True),
]
txt_multi(s4, res_lines, Inches(9.1), Inches(2.05), Inches(3.7), Inches(1.9), spacing=Pt(2))

# Right bottom: Sci-Fi / Bio Connection
rounded_box(s4, Inches(8.9), Inches(4.2), Inches(4.1), Inches(2.95), fill=CARD_BG, line=TEAL)
txt(s4, "🧬  PHM & Biological Analogy",
    Inches(9.1), Inches(4.3), Inches(3.7), Inches(0.4),
    size=Pt(16), bold=True, color=TEAL)

phm_lines = [
    ("🌌  Astrophage Migration Path",                   Pt(12), True,  WHITE,      False),
    ("    Sun-dimming Astrophage clusters must flock past dark cosmic dust blockages to reach planetary heat targets.", Pt(11), False, LIGHT_GRAY, False),
    ("",                                                 Pt(2),  False, WHITE,      False),
    ("🔬  Biomedical Target Delivery",                  Pt(12), True,  WHITE,      False),
    ("    Mirroring E. coli quorum-sensing, light-steered nanoparticles can collectively flow around tumor shadow zones.", Pt(11), False, LIGHT_GRAY, False),
]
txt_multi(s4, phm_lines, Inches(9.1), Inches(4.75), Inches(3.7), Inches(2.3), spacing=Pt(1.5))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LLM PROFICIENCY
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide()
slide_number(s5, 5)

box(s5, 0, 0, Inches(0.18), H, fill=GOLD)

txt(s5, "AI CO-PILOT: PROMPT ENGINEERING CRITICAL GAINS",
    Inches(0.4), Inches(0.22), Inches(10), Inches(0.65),
    size=Pt(28), bold=True, color=WHITE)
txt(s5, "Ryland Grace's Lab assistant: how highly detailed prompts unlocked state predictions",
    Inches(0.4), Inches(0.88), Inches(10), Inches(0.4),
    size=Pt(16), italic=True, color=GOLD)

accent_line(s5, Inches(0.4), Inches(1.36), Inches(12.5), color=GOLD)

# Left Prompt Card
rounded_box(s5, Inches(0.35), Inches(1.55), Inches(6.15), Inches(4.15), fill=CARD_BG, line=GOLD, line_width=Pt(2.5))
txt(s5, "✨  PROMPT 1 — Physics-Grounded Critic",
    Inches(0.55), Inches(1.65), Inches(5.8), Inches(0.45),
    size=Pt(16), bold=True, color=GOLD)
accent_line(s5, Inches(0.55), Inches(2.15), Inches(5.8), color=GOLD, height=Pt(1.5))

p1_text = (
    '"Review my Langevin-dynamics simulation as a PhD physicist. '
    'Verify if my self-propulsion and Vicsek alignment equations are physically sound. '
    'Correct any phenomenological simplifications using exact phototactic force equations '
    'derived from radiation heat fluxes, and provide a strict boundary mapping '
    'for reflective Neumann walls."'
)
txt(s5, p1_text, Inches(0.55), Inches(2.25), Inches(5.8), Inches(1.6), size=Pt(11.5), italic=True, color=WHITE)

txt(s5, "💡  Why it made a massive difference:",
    Inches(0.55), Inches(3.95), Inches(5.8), Inches(0.35), size=Pt(13), bold=True, color=TEAL)
p1_why = [
    ("→  Flagged numerical drift in angular noise Euler-Maruyama step", Pt(11.5), False, LIGHT_GRAY, False),
    ("→  Mapped complex heat-flux parameters directly to local array limits", Pt(11.5), False, LIGHT_GRAY, False),
    ("→  Helped write vectorized local alignment loops reducing runtimes", Pt(11.5), False, LIGHT_GRAY, False),
]
txt_multi(s5, p1_why, Inches(0.55), Inches(4.3), Inches(5.8), Inches(1.2), spacing=Pt(2))

# Right Prompt Card
rounded_box(s5, Inches(6.8), Inches(1.55), Inches(6.15), Inches(4.15), fill=CARD_BG, line=TEAL, line_width=Pt(2.5))
txt(s5, "🚀  PROMPT 2 — Emergent Phase Explorer",
    Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.45),
    size=Pt(16), bold=True, color=TEAL)
accent_line(s5, Inches(7.0), Inches(2.15), Inches(5.8), color=TEAL, height=Pt(1.5))

p2_text = (
    '"Analyze this active-matter system. If I add a constant turning torque '
    'resulting in chiral circle-swimming, and swap periodic boundaries for Neumann reflective walls, '
    'what novel macroscopic state will emerge? Detail the order parameters '
    'needed to map the transition in a new phase space sweep."'
)
txt(s5, p2_text, Inches(7.0), Inches(2.25), Inches(5.8), Inches(1.6), size=Pt(11.5), italic=True, color=WHITE)

txt(s5, "💡  Why it made a massive difference:",
    Inches(7.0), Inches(3.95), Inches(5.8), Inches(0.35), size=Pt(13), bold=True, color=GOLD)
p2_why = [
    ("→  Correctly predicted 'chiral orbital locking' around obstacles", Pt(11.5), False, LIGHT_GRAY, False),
    ("→  Suggested polar order vs vorticity metrics to trace transitions", Pt(11.5), False, LIGHT_GRAY, False),
    ("→  Designed a multi-dimensional phase diagram we then implemented", Pt(11.5), False, LIGHT_GRAY, False),
]
txt_multi(s5, p2_why, Inches(7.0), Inches(4.3), Inches(5.8), Inches(1.2), spacing=Pt(2))

# Bottom Banner
rounded_box(s5, Inches(0.35), Inches(5.85), Inches(12.6), Inches(1.35), fill=RGBColor(0x13, 0x1F, 0x38), line=GOLD, line_width=Pt(1.5))
txt(s5, "🏆  Research Takeaway", Inches(0.6), Inches(5.95), Inches(2.0), Inches(0.4), size=Pt(14), bold=True, color=GOLD)
txt(s5,
    "Highly context-rich prompts transformed the AI from a standard code assistant into a true "
    "theoretical collaborator. Feeding the LLM exact paper concepts enabled rapid conceptual sweeps, "
    "automatic boundary logic drafting, and rigorous phase mapping — speeding up research to hackathon velocity.",
    Inches(2.8), Inches(5.97), Inches(9.8), Inches(1.1), size=Pt(12), color=WHITE, wrap=True)

# Save
prs.save(OUT)
print(f"✅  Saved: {OUT}")
